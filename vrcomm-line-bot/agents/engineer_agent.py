"""
agents/engineer_agent.py — VRCOMM Engineer / Pre-Sales Agent

Handles:
  1. technical_qa     — Technical questions about VRCOMM products
  2. tor_analysis     — TOR/SOW compliance: match requirements → products
  3. spec_comparison  — Compare products against each other or a requirement

Data sources (priority order):
  1. specs/{Brand Name}/   — local brand folder with .txt/.md files
  2. SharePoint            — specs/{Brand Name}/ folder on SharePoint (Graph API)
  3. product URLs          — fetch vendor website as fallback

Three-layer anti-hallucination (same pattern as product_agent):
  Step 1: Haiku selects relevant brands from VRCOMM list only
  Step 2: Sonnet answers using positive allowlist prompt (only named brands)
  Step 3: Post-processing — forbidden-brand scan → retry → sentence strip

SharePoint config (environment variables):
  SHAREPOINT_SITE_ID     — Graph API site ID (visit /setup-sharepoint to find it)
  SHAREPOINT_SPECS_PATH  — folder path inside Documents, e.g. "ProductSpecs"
  (Uses same MS_CLIENT_ID / MS_CLIENT_SECRET / MS_TENANT_ID as email_handler)
"""
import os, re, logging, time
import requests
from anthropic import Anthropic

logger = logging.getLogger(__name__)
client = Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY", ""))

_BASE_DIR     = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
_PRODUCT_LIST = os.path.join(_BASE_DIR, "product", "VRCOMM_ProductList.xlsx")
_SPECS_DIR    = os.path.join(_BASE_DIR, "specs")

# ── Product list cache ────────────────────────────────────────────────────────

_product_cache      = None
_product_cache_time = 0.0
_CACHE_TTL          = 300


def _load_product_list() -> list:
    """Load VRCOMM_ProductList.xlsx → list of {brand, url}. Cached 5 min."""
    global _product_cache, _product_cache_time

    path = _PRODUCT_LIST
    if not os.path.isfile(path):
        logger.error("ProductList not found: %s", path)
        return []

    mtime = os.path.getmtime(path)
    now   = time.time()
    if (_product_cache is not None
            and mtime == _product_cache_time
            and (now - _product_cache_time) < _CACHE_TTL):
        return _product_cache

    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        ws = wb.active
        products = []
        for row in ws.iter_rows(values_only=True):
            brand = str(row[0]).strip() if row[0] else ""
            url   = str(row[1]).strip() if len(row) > 1 and row[1] else ""
            if brand.lower() in ("brand", "product", "name", "") or not brand:
                continue
            products.append({"brand": brand, "url": url})
        _product_cache      = products
        _product_cache_time = mtime
        logger.info("[engineer_agent] ProductList: %d brands", len(products))
        return products
    except Exception as e:
        logger.error("[engineer_agent] ProductList load error: %s", e)
        return []


# ── PDF text extractor ───────────────────────────────────────────────────────

def _extract_pdf_text(path_or_bytes, max_chars: int = 4000) -> str:
    """
    Extract plain text from a PDF file or bytes object using pdfplumber.
    Falls back to empty string if pdfplumber is not installed or extraction fails.

    Args:
        path_or_bytes : file path (str) OR raw bytes (from SharePoint download)
        max_chars     : maximum characters to return (keeps context window manageable)
    """
    try:
        import pdfplumber, io
        source = io.BytesIO(path_or_bytes) if isinstance(path_or_bytes, bytes) else path_or_bytes
        with pdfplumber.open(source) as pdf:
            pages_text = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text.strip())
                # Also extract tables as plain text rows
                for table in page.extract_tables():
                    for row in table:
                        row_clean = [str(c).strip() if c else "" for c in row]
                        pages_text.append(" | ".join(row_clean))
            return "\n\n".join(pages_text)[:max_chars]
    except ImportError:
        logger.warning("[engineer_agent] pdfplumber not installed — PDF skipped")
        return ""
    except Exception as e:
        logger.warning("[engineer_agent] PDF extract error: %s", e)
        return ""


def _extract_docx_text(path_or_bytes, max_chars: int = 4000) -> str:
    """Extract plain text from a .docx file or bytes using python-docx."""
    try:
        from docx import Document
        import io
        source = io.BytesIO(path_or_bytes) if isinstance(path_or_bytes, bytes) else path_or_bytes
        doc = Document(source)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)[:max_chars]
    except ImportError:
        logger.warning("[engineer_agent] python-docx not installed — .docx skipped")
        return ""
    except Exception as e:
        logger.warning("[engineer_agent] DOCX extract error: %s", e)
        return ""


def _extract_pptx_text(path_or_bytes, max_chars: int = 4000) -> str:
    """Extract plain text from a .pptx file or bytes using python-pptx."""
    try:
        from pptx import Presentation
        import io
        source = io.BytesIO(path_or_bytes) if isinstance(path_or_bytes, bytes) else path_or_bytes
        prs = Presentation(source)
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                # Extract table cells
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_texts.append(row_text)
            # Include slide notes if present
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append("[Notes: %s]" % notes)
            if slide_texts:
                parts.append("[Slide %d]\n%s" % (slide_num, "\n".join(slide_texts)))
        return "\n\n".join(parts)[:max_chars]
    except ImportError:
        logger.warning("[engineer_agent] python-pptx not installed — .pptx skipped")
        return ""
    except Exception as e:
        logger.warning("[engineer_agent] PPTX extract error: %s", e)
        return ""


def _extract_file(fpath_or_bytes, filename: str, max_chars: int = 3000) -> str:
    """Route to the correct extractor based on file extension."""
    ext = filename.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return _extract_pdf_text(fpath_or_bytes, max_chars)
    elif ext == "docx":
        return _extract_docx_text(fpath_or_bytes, max_chars)
    elif ext == "pptx":
        return _extract_pptx_text(fpath_or_bytes, max_chars)
    return ""


# ── Local spec loader (folder-based) ─────────────────────────────────────────

_spec_cache: dict = {}   # brand_key → {"content": str, "mtime": float}


def _brand_folder_match(brand: str) -> str:
    """
    Find the best-matching subfolder in specs/ for a given brand name.
    Returns the full folder path, or empty string if not found.

    Matching strategy:
      1. Exact match (case-insensitive)
      2. Brand name is a substring of folder name
      3. Folder name is a substring of brand name
    """
    if not os.path.isdir(_SPECS_DIR):
        return ""

    brand_lower = brand.lower()
    try:
        for entry in os.scandir(_SPECS_DIR):
            if not entry.is_dir():
                continue
            folder_lower = entry.name.lower()
            if (folder_lower == brand_lower
                    or brand_lower in folder_lower
                    or folder_lower in brand_lower):
                return entry.path
    except Exception as e:
        logger.warning("[engineer_agent] Folder scan error: %s", e)
    return ""


def _load_spec_local(brand: str) -> str:
    """
    Load all .txt and .md files from specs/{Brand Name}/ folder.
    Falls back to legacy flat-file lookup (specs/brand_name.txt) for compatibility.
    Content is cached by folder mtime.
    """
    brand_key = brand.lower()

    # ── Priority 1: brand subfolder (new structure) ───────────────────────────
    folder = _brand_folder_match(brand)
    if folder:
        try:
            _TEXT_EXTS = (".txt", ".md")
            _RICH_EXTS = (".pdf", ".docx", ".pptx")

            all_files   = os.listdir(folder)
            txt_files   = sorted([f for f in all_files if f.lower().endswith(_TEXT_EXTS)])
            rich_files  = sorted([f for f in all_files if f.lower().endswith(_RICH_EXTS)])
            all_spec    = txt_files + rich_files

            if all_spec:
                all_paths  = [os.path.join(folder, f) for f in all_spec]
                max_mtime  = max(os.path.getmtime(p) for p in all_paths)
                cached     = _spec_cache.get(brand_key)
                if cached and cached.get("mtime") == max_mtime:
                    return cached["content"]

                parts = []
                for fname in txt_files:
                    fpath = os.path.join(folder, fname)
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as fh:
                        parts.append(fh.read())

                for fname in rich_files:
                    fpath = os.path.join(folder, fname)
                    extracted = _extract_file(fpath, fname, max_chars=3000)
                    if extracted:
                        ext = fname.rsplit(".", 1)[-1].upper()
                        parts.append("[%s: %s]\n%s" % (ext, fname, extracted))
                        logger.info("[engineer_agent] Extracted %s: %s", ext, fname)

                content = "\n\n".join(parts)[:10000]
                _spec_cache[brand_key] = {"content": content, "mtime": max_mtime}
                logger.info("[engineer_agent] Loaded spec folder: %s (%d txt, %d rich)",
                            folder, len(txt_files), len(rich_files))
                return content
        except Exception as e:
            logger.warning("[engineer_agent] Spec folder load error [%s]: %s", folder, e)

    # ── Priority 2: legacy flat file (specs/brand_name.txt) ──────────────────
    brand_slug = brand.lower().replace(" ", "_").replace("-", "_")
    try:
        for fname in os.listdir(_SPECS_DIR):
            if not fname.endswith((".txt", ".md")):
                continue
            fname_key = fname.lower().replace("-", "_").rsplit(".", 1)[0]
            if fname_key == brand_slug or brand_slug in fname_key or fname_key in brand_slug:
                fpath = os.path.join(_SPECS_DIR, fname)
                mtime = os.path.getmtime(fpath)
                cached = _spec_cache.get(brand_key)
                if cached and cached.get("mtime") == mtime:
                    return cached["content"]
                with open(fpath, "r", encoding="utf-8", errors="ignore") as fh:
                    content = fh.read()[:5000]
                _spec_cache[brand_key] = {"content": content, "mtime": mtime}
                logger.info("[engineer_agent] Loaded spec file (legacy): %s", fname)
                return content
    except Exception as e:
        logger.warning("[engineer_agent] Legacy spec file error: %s", e)

    return ""


# ── SharePoint spec loader ────────────────────────────────────────────────────

_SHAREPOINT_SITE_ID   = os.environ.get("SHAREPOINT_SITE_ID", "")
_SHAREPOINT_SPECS_PATH = os.environ.get("SHAREPOINT_SPECS_PATH", "ProductSpecs")
_sp_token_cache: dict = {}   # {"token": str, "expires_at": float}
_sp_spec_cache:  dict = {}   # brand_key → {"content": str, "fetched_at": float}
_SP_CACHE_TTL = 3600         # 1 hour


def _get_graph_token() -> str:
    """Get Microsoft Graph API access token (cached). Same credentials as email_handler."""
    now = time.time()
    if _sp_token_cache.get("token") and now < _sp_token_cache.get("expires_at", 0) - 60:
        return _sp_token_cache["token"]

    # Use same env var names as email_handler.py
    tenant_id     = os.environ.get("MS_TENANT_ID", "")
    client_id     = os.environ.get("MS_CLIENT_ID", "")
    client_secret = os.environ.get("MS_CLIENT_SECRET", "")

    if not all([tenant_id, client_id, client_secret]):
        return ""

    try:
        resp = requests.post(
            "https://login.microsoftonline.com/%s/oauth2/v2.0/token" % tenant_id,
            data={
                "grant_type":    "client_credentials",
                "client_id":     client_id,
                "client_secret": client_secret,
                "scope":         "https://graph.microsoft.com/.default",
            },
            timeout=10,
        )
        resp.raise_for_status()
        data = resp.json()
        token = data.get("access_token", "")
        _sp_token_cache["token"]      = token
        _sp_token_cache["expires_at"] = now + data.get("expires_in", 3600)
        return token
    except Exception as e:
        logger.warning("[engineer_agent] Graph token error: %s", e)
        return ""


def _load_spec_sharepoint(brand: str) -> str:
    """
    Load spec files from SharePoint: Documents/{SHAREPOINT_SPECS_PATH}/{Brand Name}/
    Returns concatenated text content of all .txt and .md files in that folder.
    Cached 1 hour.
    """
    if not _SHAREPOINT_SITE_ID:
        return ""

    brand_key = brand.lower()
    now = time.time()
    cached = _sp_spec_cache.get(brand_key)
    if cached and (now - cached["fetched_at"]) < _SP_CACHE_TTL:
        return cached["content"]

    token = _get_graph_token()
    if not token:
        return ""

    headers = {"Authorization": "Bearer " + token, "Accept": "application/json"}
    base_url = "https://graph.microsoft.com/v1.0"

    # Step 1: Find the brand subfolder under the specs path
    # Search in /sites/{site_id}/drive/root:/{specs_path}/{brand}/children
    brand_folder_path = "%s/%s" % (_SHAREPOINT_SPECS_PATH.strip("/"), brand)
    folder_url = "%s/sites/%s/drive/root:/%s:/children" % (
        base_url, _SHAREPOINT_SITE_ID, brand_folder_path
    )

    try:
        resp = requests.get(folder_url, headers=headers, timeout=10)
        if resp.status_code == 404:
            # Try fuzzy match: list parent folder and find best match
            parent_url = "%s/sites/%s/drive/root:/%s:/children" % (
                base_url, _SHAREPOINT_SITE_ID, _SHAREPOINT_SPECS_PATH
            )
            parent_resp = requests.get(parent_url, headers=headers, timeout=10)
            if parent_resp.status_code != 200:
                return ""

            items = parent_resp.json().get("value", [])
            brand_lower = brand.lower()
            match_id = None
            for item in items:
                if item.get("folder") is None:
                    continue
                name_lower = item["name"].lower()
                if (name_lower == brand_lower
                        or brand_lower in name_lower
                        or name_lower in brand_lower):
                    match_id = item["id"]
                    break

            if not match_id:
                logger.info("[engineer_agent] SharePoint: no folder found for '%s'", brand)
                _sp_spec_cache[brand_key] = {"content": "", "fetched_at": now}
                return ""

            children_url = "%s/sites/%s/drive/items/%s/children" % (
                base_url, _SHAREPOINT_SITE_ID, match_id
            )
            resp = requests.get(children_url, headers=headers, timeout=10)

        resp.raise_for_status()
        files = resp.json().get("value", [])

    except Exception as e:
        logger.warning("[engineer_agent] SharePoint folder list error [%s]: %s", brand, e)
        return ""

    # Step 2: Download each supported spec file
    _SUPPORTED = (".txt", ".md", ".pdf", ".docx", ".pptx")
    parts = []
    for item in files:
        name = item.get("name", "")
        if not name.lower().endswith(_SUPPORTED):
            continue
        download_url = item.get("@microsoft.graph.downloadUrl")
        if not download_url:
            continue
        try:
            dl = requests.get(download_url, timeout=15)
            dl.raise_for_status()
            if name.lower().endswith((".txt", ".md")):
                parts.append(dl.text[:3000])
                logger.info("[engineer_agent] SharePoint downloaded: %s for '%s'", name, brand)
            else:
                extracted = _extract_file(dl.content, name, max_chars=3000)
                if extracted:
                    ext = name.rsplit(".", 1)[-1].upper()
                    parts.append("[%s: %s]\n%s" % (ext, name, extracted))
                    logger.info("[engineer_agent] SharePoint %s extracted: %s for '%s'",
                                ext, name, brand)
        except Exception as e:
            logger.warning("[engineer_agent] SharePoint download error [%s]: %s", name, e)

    content = "\n\n".join(parts)[:6000] if parts else ""
    _sp_spec_cache[brand_key] = {"content": content, "fetched_at": now}

    if content:
        logger.info("[engineer_agent] SharePoint spec loaded for '%s' (%d chars)", brand, len(content))
    return content


def _load_spec_file(brand: str) -> str:
    """
    Load spec content for a brand. Priority:
      1. Local specs/{Brand Name}/ folder
      2. SharePoint specs folder (if SHAREPOINT_SITE_ID is set)
      3. (Caller falls back to URL fetch if this returns empty)
    """
    content = _load_spec_local(brand)
    if content:
        return content

    content = _load_spec_sharepoint(brand)
    return content


# ── URL fetcher (cached 1 hr) ─────────────────────────────────────────────────

_url_cache: dict = {}
_URL_CACHE_TTL   = 3600


def _fetch_url(url: str) -> str:
    """Fetch URL → clean plain text (HTML stripped). Cached 1 hour."""
    url = url.strip()
    if not url:
        return ""

    now    = time.time()
    cached = _url_cache.get(url)
    if cached and (now - cached["fetched_at"]) < _URL_CACHE_TTL:
        return cached["content"]

    try:
        resp = requests.get(url, headers={
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                "AppleWebKit/537.36 Chrome/120.0.0.0 Safari/537.36"
            ),
            "Accept-Language": "en-US,en;q=0.9,th;q=0.8",
        }, timeout=8, allow_redirects=True)
        resp.raise_for_status()

        html = resp.text
        html = re.sub(r'<(script|style|nav|footer|header)[^>]*>.*?</\1>',
                      ' ', html, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r'<[^>]+>', ' ', html)
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text).strip()[:2500]

        _url_cache[url] = {"content": text, "fetched_at": now}
        logger.info("[engineer_agent] Fetched %d chars: %s", len(text), url[:60])
        return text
    except Exception as e:
        logger.warning("[engineer_agent] URL fetch failed [%s]: %s", url[:60], e)
        return ""


# ── Mode detector ─────────────────────────────────────────────────────────────

_TOR_KEYWORDS = [
    "tor", "sow", "scope of work", "terms of reference",
    "tdr", "spec compliance", "ตาราง", "comply", "requirement",
    "ข้อกำหนด", "คุณสมบัติ", "เทียบ spec", "compliance",
    "tender", "bidding", "จัดซื้อ", "คุณลักษณะ", "รายการ",
]


def _detect_mode(message: str) -> str:
    """
    Detect whether this is a TOR analysis or general technical Q&A.
    Returns: "tor_analysis" or "technical_qa"
    """
    msg_lower = message.lower()
    if any(kw in msg_lower for kw in _TOR_KEYWORDS):
        return "tor_analysis"
    # Long messages with numbered items are likely TOR/requirements
    if len(message) > 400 and re.search(r'^\s*\d+[\.\)]\s+', message, re.MULTILINE):
        return "tor_analysis"
    return "technical_qa"


# ── STEP 1: Brand selector ────────────────────────────────────────────────────

_SELECT_PROMPT = """You are a technical product selector for VRCOMM, a cybersecurity company in Thailand.

Customer/engineer query:
\"{message}\"

VRCOMM's complete product list (these are the ONLY brands we sell):
{brand_list}

Select which brands from the list are technically relevant to this query.
Consider: product category, use case, technical requirements, complementary products.
Return ONLY exact brand names from the list, one per line.
If none are relevant, return: NONE
Do NOT add brands outside the list. No explanations."""


def _select_relevant_brands(message: str, product_list: list) -> list:
    """Step 1: Claude picks relevant brands from our list only."""
    brand_list_text = "\n".join(
        "%d. %s" % (i + 1, p["brand"]) for i, p in enumerate(product_list)
    )
    try:
        resp = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=200,
            messages=[{"role": "user", "content": _SELECT_PROMPT.format(
                message=message[:800],
                brand_list=brand_list_text,
            )}],
        )
        raw = resp.content[0].text.strip()
        logger.info("[engineer_agent] Step1 selection: %s", raw[:120])

        if raw.strip().upper() == "NONE" or not raw:
            return []

        selected_names = [
            re.sub(r'^[\d\.\-\)\s]+', '', line).strip()
            for line in raw.splitlines()
            if line.strip() and line.strip().upper() != "NONE"
        ]
        matched = []
        for name in selected_names:
            for p in product_list:
                if (p["brand"].lower() == name.lower()
                        or name.lower() in p["brand"].lower()
                        or p["brand"].lower() in name.lower()):
                    if p not in matched:
                        matched.append(p)
        logger.info("[engineer_agent] Brands selected: %s",
                    [m["brand"] for m in matched])
        return matched
    except Exception as e:
        logger.error("[engineer_agent] Step1 error: %s", e)
        return []


# ── Forbidden brands (same list as product_agent) ────────────────────────────

_FORBIDDEN_BRANDS = [
    "fortinet", "fortigate", "cisco", "palo alto", "check point", "checkpoint",
    "juniper", "sonicwall", "sonic wall", "barracuda", "trend micro", "trendmicro",
    "crowdstrike", "crowd strike", "sentinelone", "sentinel one", "mcafee",
    "symantec", "broadcom", "hp ", "hewlett", "dell ", "aruba", "meraki",
    "watchguard", "watch guard", "zyxel", "netgear", "ubiquiti",
]


def _contains_forbidden_brand(text: str) -> str:
    """Return the first forbidden brand found in text, or empty string."""
    text_lower = text.lower()
    for brand in _FORBIDDEN_BRANDS:
        if brand in text_lower:
            return brand
    return ""


def _strip_forbidden_sentences(text: str) -> str:
    """Nuclear option: remove any sentence containing a forbidden brand."""
    parts = re.split(r'(?<=[.!?\n])\s+', text)
    clean, removed = [], []
    for part in parts:
        if _contains_forbidden_brand(part):
            removed.append(part[:60])
        else:
            clean.append(part)
    if removed:
        logger.warning("[engineer_agent] Stripped %d forbidden sentence(s): %s",
                       len(removed), removed)
    result = " ".join(clean).strip()
    return result or "ขออภัยครับ ไม่สามารถให้ข้อมูลเพิ่มเติมได้ กรุณาติดต่อทีม Engineer โดยตรงครับ"


# ── STEP 2: System prompts (Sonnet — positive allowlist approach) ─────────────

_TECH_QA_SYSTEM = """You are VRCOMM's Senior Network and Cybersecurity Engineer briefing a colleague.
VRCOMM is a Network and Cybersecurity solutions provider in Thailand.
This is an INTERNAL conversation — staff to staff, NOT customer-facing.

══════════════════════════════════════════════════════
ALLOWED BRANDS — you may ONLY discuss these, nothing else:
{allowed_brand_names}

Product specs:
{product_specs_section}
══════════════════════════════════════════════════════

BEFORE writing anything, verify: every brand name in your reply must appear in the ALLOWED BRANDS list above.

RULES:
- Mention ONLY brands from the ALLOWED BRANDS list
- Do NOT mention Fortinet, FortiGate, Cisco, Palo Alto, Check Point, Juniper, SonicWall or any brand not in the list — these are our competitors
- Provide accurate technical specs and model recommendations
- For sizing questions: address users, throughput, concurrent sessions
- NEVER quote prices — say: "ให้ Sales ดึง cost sheet แล้วทำ quote ได้เลยครับ"
- Tone: direct, technical, colleague-to-colleague — internal briefing, not customer service
- Reply in the SAME LANGUAGE as the staff member (Thai → Thai, English → English)
- Plain text only — no markdown, no bullets, no headers
- Max 5 short paragraphs"""

_TOR_ANALYSIS_SYSTEM = """You are VRCOMM's Senior Pre-Sales Engineer specialising in TOR/SOW compliance.
VRCOMM is a Network and Cybersecurity solutions provider in Thailand.
This is an INTERNAL analysis for a colleague preparing a proposal.

══════════════════════════════════════════════════════
ALLOWED BRANDS — you may ONLY recommend these, nothing else:
{allowed_brand_names}

Product specs:
{product_specs_section}
══════════════════════════════════════════════════════

Analyse the TOR/SOW and produce a compliance table. Every product you recommend MUST be from the ALLOWED BRANDS list.

Output format — one line per requirement:
  ข้อ [No.] | ข้อกำหนด: [requirement text] | สินค้า: [Brand + model] | สถานะ: [FULL/PARTIAL/ALT] | หมายเหตุ: [brief note]

Status definitions:
  FULL    — meets requirement fully out of the box
  PARTIAL — meets with additional config or add-on license
  ALT     — alternative solution using VRCOMM product covers the intent

After the table: 2-3 sentences summarising the proposed solution and any gaps.

RULES:
- Recommend ONLY brands from the ALLOWED BRANDS list
- Do NOT mention Fortinet, Cisco, Palo Alto, Check Point, or any brand outside the list
- NEVER quote prices
- Reply in the SAME LANGUAGE as the document (Thai → Thai, English → English)
- Plain text only — no markdown formatting"""

_NO_MATCH_SYSTEM = """You are VRCOMM's Senior Network and Cybersecurity Engineer briefing a colleague.
VRCOMM is a Network and Cybersecurity solutions provider in Thailand.
This is an INTERNAL conversation — staff to staff, NOT customer-facing.

The query is about a product type where no direct brand match was found in VRCOMM's catalog.
Recommend the best available alternatives.

══════════════════════════════════════════════════════
VRCOMM's COMPLETE product portfolio (ALL brands we sell — nothing else):
{full_list}
══════════════════════════════════════════════════════

BEFORE writing anything, verify: every brand name you mention must appear in the list above.

RULES:
- Recommend ONLY brands from the list above
- Do NOT mention Fortinet, Cisco, Palo Alto, Check Point, or any brand not on the list
- NEVER quote prices — say: "ให้ Sales ดึง cost sheet แล้วทำ quote ได้เลยครับ"
- Tone: direct, technical, colleague-to-colleague
- Reply in the SAME LANGUAGE as the staff member
- Plain text only — max 5 short paragraphs"""


def _build_product_specs_section(selected: list) -> str:
    """Build spec content for each selected brand."""
    sections = []
    for p in selected:
        # Priority: local spec file > URL fetch > just name+URL
        spec_content = _load_spec_file(p["brand"])
        if not spec_content:
            spec_content = _fetch_url(p["url"])

        section = "=== %s ===\nWebsite: %s" % (p["brand"], p["url"])
        if spec_content:
            # Allow up to 6000 chars per brand so PDF/PPTX/DOCX content isn't cut off
            section += "\n\n%s" % spec_content[:6000]
        sections.append(section)
    return "\n\n".join(sections) if sections else "(No spec content available)"


def _build_system(mode: str, selected: list, product_list: list) -> str:
    if not selected:
        full_list = "\n".join(
            "%d. %s" % (i + 1, p["brand"]) for i, p in enumerate(product_list)
        )
        return _NO_MATCH_SYSTEM.format(full_list=full_list)

    # Positive allowlist — tell Claude exactly which brands it CAN use
    allowed_brand_names = "\n".join(
        "  - %s" % p["brand"] for p in selected
    )
    product_specs_section = _build_product_specs_section(selected)

    if mode == "tor_analysis":
        return _TOR_ANALYSIS_SYSTEM.format(
            allowed_brand_names=allowed_brand_names,
            product_specs_section=product_specs_section,
        )
    else:
        return _TECH_QA_SYSTEM.format(
            allowed_brand_names=allowed_brand_names,
            product_specs_section=product_specs_section,
        )


# ── Main handler ──────────────────────────────────────────────────────────────

def handle(message: str, user_name: str, user_id: str,
           source: str = "line", history: list = None,
           intent: str = "technical", **kwargs) -> str:
    """
    Two-step engineer handler:
      Step 0 — Detect mode (TOR analysis vs technical Q&A)
      Step 1 — SELECT: Claude picks relevant brands from our list
      Step 2 — ANSWER: Claude answers with spec context
    """
    if history is None:
        history = []

    product_list = _load_product_list()
    if not product_list:
        from agents.general_agent import handle as general_handle
        return general_handle(message=message, user_name=user_name,
                              user_id=user_id, source=source,
                              history=history, intent=intent)

    # Step 0: detect mode
    mode = _detect_mode(message)
    logger.info("[engineer_agent] mode=%s for: %s", mode, message[:60])

    # Step 1: select relevant brands
    selected = _select_relevant_brands(message, product_list)

    # Step 2: build system + answer
    system = _build_system(mode, selected, product_list)

    content = message if history else (
        "Engineer/Staff: %s%s\n\n%s" % (
            user_name,
            " (via email)" if source == "email" else "",
            message,
        )
    )
    messages = history + [{"role": "user", "content": content}]

    # TOR analysis needs more tokens for the compliance table
    max_tokens = 1500 if mode == "tor_analysis" else 768

    try:
        # Step 2 uses Sonnet — better instruction-following under strict constraints
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=max_tokens,
            system=system,
            messages=messages,
        )
        reply = response.content[0].text.strip()

        # ── Layer 2: Retry if forbidden brand detected ────────────────────────
        forbidden = _contains_forbidden_brand(reply)
        if forbidden:
            logger.warning("[engineer_agent] Forbidden brand '%s' in reply — retrying", forbidden)
            retry_messages = messages + [
                {"role": "assistant", "content": reply},
                {"role": "user", "content": (
                    "คำตอบของคุณเมนชั่น '%s' ซึ่งไม่มีในรายการสินค้า VRCOMM เลย "
                    "กรุณาตอบใหม่โดยใช้เฉพาะสินค้าใน ALLOWED BRANDS list เท่านั้น "
                    "ห้ามเมนชั่น %s หรือสินค้าอื่นที่ไม่อยู่ใน list" % (forbidden, forbidden)
                )},
            ]
            retry_resp = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=max_tokens,
                system=system,
                messages=retry_messages,
            )
            reply = retry_resp.content[0].text.strip()
            logger.info("[engineer_agent] Retry reply: %s", reply[:80])

            # ── Layer 3: Sentence-level strip — nuclear last resort ────────
            if _contains_forbidden_brand(reply):
                logger.warning("[engineer_agent] Retry still has forbidden brand — stripping")
                reply = _strip_forbidden_sentences(reply)

        logger.info("[engineer_agent] replied (mode=%s, brands=%s): %s",
                    mode, [s["brand"] for s in selected], reply[:80])
        return reply

    except Exception as e:
        logger.error("[engineer_agent] Claude API error: %s", e)
        return ("ขออภัยครับ ระบบขัดข้องชั่วคราว "
                "กรุณาติดต่อทีม VRCOMM ที่ iddhi.t@vrcomm.net ครับ")
